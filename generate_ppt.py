import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION, XL_LABEL_POSITION
from pptx.enum.text import MSO_AUTO_SIZE

# Define Theme Colors based on CSS
COLOR_DARK_BLUE = RGBColor(13, 27, 42)
COLOR_PRIMARY_BLUE = RGBColor(24, 95, 165)
COLOR_LIGHT_GRAY = RGBColor(244, 246, 249)
COLOR_WHITE = RGBColor(255, 255, 255)
COLOR_MUTED_TEXT = RGBColor(143, 163, 177)
COLOR_DARK_TEXT = RGBColor(26, 26, 46)

prs = Presentation()
prs.slide_width = Inches(10)
prs.slide_height = Inches(5.625)
blank_slide_layout = prs.slide_layouts[6]

def hex_to_rgb(hex_str):
    h = hex_str.lstrip('#')
    return RGBColor(int(h[0:2], 16), int(h[2:4], 16), int(h[4:6], 16))

def add_background(slide, color):
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), prs.slide_width, prs.slide_height)
    bg.fill.solid()
    bg.fill.fore_color.rgb = color
    bg.line.fill.background()
    return bg

def apply_font(paragraph, size, color, bold=False, name="Segoe UI"):
    paragraph.font.size = Pt(size)
    if color:
        paragraph.font.color.rgb = color
    paragraph.font.bold = bold
    paragraph.font.name = name

def add_text(slide, text, left, top, width, height, font_size, font_color, bold=False, align=None, v_align=None, auto_size=False):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.clear()
    tf.word_wrap = True
    if auto_size:
        tf.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
    if v_align:
        tf.vertical_anchor = v_align
    
    p = tf.paragraphs[0]
    p.text = text
    apply_font(p, font_size, font_color, bold)
    if align:
        p.alignment = align
    return txBox

def add_footer(slide, is_light=False):
    f_bg_color = COLOR_PRIMARY_BLUE if is_light else RGBColor(12, 68, 124)
    f_text_color = COLOR_WHITE
    ftr = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), prs.slide_height - Inches(0.4), prs.slide_width, Inches(0.4))
    ftr.fill.solid()
    ftr.fill.fore_color.rgb = f_bg_color
    ftr.line.fill.background()
    add_text(slide, "SGC Growth Intelligence Toolkit", Inches(0.4), prs.slide_height - Inches(0.35), Inches(4), Inches(0.3), 11, f_text_color)
    add_text(slide, "Northeastern University · ALY6980", Inches(6.5), prs.slide_height - Inches(0.35), Inches(3), Inches(0.3), 11, f_text_color, align=PP_ALIGN.RIGHT)

def add_header(slide, num_str, title, subtitle, is_light=False):
    h_bg_color = COLOR_PRIMARY_BLUE
    h_text_color = COLOR_WHITE
    b_color = COLOR_LIGHT_GRAY if is_light else COLOR_DARK_BLUE
    
    hdr = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), prs.slide_width, Inches(1.2))
    hdr.fill.solid()
    hdr.fill.fore_color.rgb = h_bg_color
    hdr.line.fill.background()
    
    add_text(slide, num_str, Inches(0.4), Inches(0.15), Inches(2), Inches(0.2), 11, RGBColor(200, 200, 200), bold=True)
    add_text(slide, title, Inches(0.4), Inches(0.35), Inches(9.2), Inches(0.4), 24, h_text_color, bold=True)
    add_text(slide, subtitle, Inches(0.4), Inches(0.8), Inches(9.2), Inches(0.3), 13, RGBColor(220, 220, 220))

def add_highlight_box(slide, left, top, width, height, text, is_light=False):
    fill_color = RGBColor(232, 240, 251) if is_light else RGBColor(24, 45, 75)
    text_color = COLOR_DARK_TEXT if is_light else RGBColor(200, 214, 232)
    
    box = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    box.fill.solid()
    box.fill.fore_color.rgb = fill_color
    box.line.color.rgb = COLOR_PRIMARY_BLUE
    box.line.width = Pt(4)
    
    txBox = slide.shapes.add_textbox(left + Inches(0.1), top + Inches(0.1), width - Inches(0.2), height - Inches(0.2))
    tf = txBox.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    tf.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
    tf.clear()
    p = tf.paragraphs[0]
    p.text = text
    apply_font(p, 12.5, text_color)
    return box

def add_stat_box(slide, left, top, width, height, num, lbl, bg_color=None):
    if not bg_color:
        bg_color = COLOR_PRIMARY_BLUE
    box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    box.fill.solid()
    box.fill.fore_color.rgb = bg_color
    box.line.fill.background()
    
    add_text(slide, num, left, top + Inches(0.05), width, height/2.2, 34, COLOR_WHITE, bold=True, align=PP_ALIGN.CENTER, v_align=MSO_ANCHOR.BOTTOM)
    add_text(slide, lbl, left + Inches(0.1), top + height/1.8, width - Inches(0.2), height/2.2, 10, RGBColor(220, 220, 220), align=PP_ALIGN.CENTER, v_align=MSO_ANCHOR.TOP, auto_size=True)

def add_card(slide, left, top, width, height, title, content_lines, is_light=False):
    bg_color = COLOR_WHITE if is_light else RGBColor(18, 30, 48)
    txt_color = COLOR_DARK_TEXT if is_light else RGBColor(200, 214, 232)
    title_color = COLOR_PRIMARY_BLUE if is_light else RGBColor(91, 163, 232)
    
    box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    box.fill.solid()
    box.fill.fore_color.rgb = bg_color
    box.line.color.rgb = RGBColor(200, 200, 200) if is_light else RGBColor(40, 60, 80)
    box.line.width = Pt(1)
    
    add_text(slide, title, left + Inches(0.15), top + Inches(0.1), width - Inches(0.3), Inches(0.35), 11.5, title_color, bold=True)
    
    txBox = slide.shapes.add_textbox(left + Inches(0.15), top + Inches(0.4), width - Inches(0.3), height - Inches(0.5))
    tf = txBox.text_frame
    tf.word_wrap = True
    tf.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
    tf.clear()
    
    for i, line in enumerate(content_lines):
        p = tf.add_paragraph() if i > 0 else tf.paragraphs[0]
        p.text = "→ " + line
        apply_font(p, 11, txt_color)
        p.space_after = Pt(4)

def clean_chart(chart):
    if chart.has_legend:
        chart.legend.font.name = 'Segoe UI'
        chart.legend.font.size = Pt(10)
        chart.legend.font.color.rgb = COLOR_DARK_TEXT
        
    try:
        category_axis = chart.category_axis
        category_axis.has_major_gridlines = False
        category_axis.tick_labels.font.name = 'Segoe UI'
        category_axis.tick_labels.font.size = Pt(9)
        category_axis.tick_labels.font.color.rgb = COLOR_MUTED_TEXT
    except ValueError:
        pass
    
    try:
        value_axis = chart.value_axis
        value_axis.has_major_gridlines = False
        value_axis.tick_labels.font.name = 'Segoe UI'
        value_axis.tick_labels.font.size = Pt(9)
        value_axis.tick_labels.font.color.rgb = COLOR_MUTED_TEXT
    except ValueError:
        pass

# ==========================================
# SLIDE 1: COVER
# ==========================================
s1 = prs.slides.add_slide(blank_slide_layout)
add_background(s1, COLOR_DARK_BLUE)

logo = s1.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(4.5), Inches(0.8), Inches(1), Inches(1))
logo.fill.solid()
logo.fill.fore_color.rgb = COLOR_PRIMARY_BLUE
logo.line.fill.background()
add_text(s1, "SGC", Inches(4.5), Inches(1.1), Inches(1), Inches(0.4), 24, COLOR_WHITE, bold=True, align=PP_ALIGN.CENTER)

add_text(s1, "ALY6980 · CAPSTONE PROJECT · WINTER 2026", Inches(1), Inches(2), Inches(8), Inches(0.3), 12, COLOR_PRIMARY_BLUE, bold=True, align=PP_ALIGN.CENTER)
add_text(s1, "SGC Growth Intelligence", Inches(1), Inches(2.3), Inches(8), Inches(0.55), 44, COLOR_WHITE, bold=True, align=PP_ALIGN.CENTER)
add_text(s1, "Toolkit & Dashboard", Inches(1), Inches(2.85), Inches(8), Inches(0.55), 44, COLOR_PRIMARY_BLUE, bold=True, align=PP_ALIGN.CENTER)

add_text(s1, "Transforming static vendor and company lists into an actionable Business Development Funnel for SGC leadership", 
         Inches(2), Inches(3.6), Inches(6), Inches(0.7), 14, COLOR_MUTED_TEXT, align=PP_ALIGN.CENTER, auto_size=True)

add_text(s1, "INSTITUTION", Inches(2), Inches(4.5), Inches(2), Inches(0.2), 10, COLOR_MUTED_TEXT, bold=True, align=PP_ALIGN.CENTER)
add_text(s1, "Northeastern University CPS", Inches(2), Inches(4.7), Inches(2), Inches(0.3), 13, COLOR_WHITE, bold=True, align=PP_ALIGN.CENTER)

add_text(s1, "COURSE", Inches(4), Inches(4.5), Inches(2), Inches(0.2), 10, COLOR_MUTED_TEXT, bold=True, align=PP_ALIGN.CENTER)
add_text(s1, "ALY6980 Capstone", Inches(4), Inches(4.7), Inches(2), Inches(0.3), 13, COLOR_WHITE, bold=True, align=PP_ALIGN.CENTER)

add_text(s1, "TERM", Inches(6), Inches(4.5), Inches(2), Inches(0.2), 10, COLOR_MUTED_TEXT, bold=True, align=PP_ALIGN.CENTER)
add_text(s1, "Winter 2026", Inches(6), Inches(4.7), Inches(2), Inches(0.3), 13, COLOR_WHITE, bold=True, align=PP_ALIGN.CENTER)

bar = s1.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), prs.slide_height - Inches(0.1), prs.slide_width, Inches(0.1))
bar.fill.solid(); bar.fill.fore_color.rgb = RGBColor(29, 158, 117); bar.line.fill.background()

# ==========================================
# SLIDE 2: PROBLEM STATEMENT
# ==========================================
s2 = prs.slides.add_slide(blank_slide_layout)
add_background(s2, COLOR_DARK_BLUE)
add_header(s2, "SLIDE 01", "Problem Statement", "The challenge SGC faces without a data-driven engagement system")
add_footer(s2)

add_highlight_box(s2, Inches(0.4), Inches(1.4), Inches(9.2), Inches(0.9), "SGC, a consulting firm seeking to bid on contracts across Massachusetts industries, currently relies on manual data lookup across thousands of vendor and company records — making strategic, resource-efficient decision-making nearly impossible.")

add_text(s2, "CORE PAIN POINTS", Inches(0.4), Inches(2.4), Inches(4), Inches(0.3), 11, COLOR_PRIMARY_BLUE, bold=True)
pain_points = [
    "No centralized view of which industry sectors hold the highest concentration of target companies for outreach",
    "Vendor landscape is fragmented — medical, logistical, and professional providers tracked in separate static lists",
    "SDO (Supplier Diversity & Opportunity) compliance data is unstructured, limiting CSR reporting capabilities",
    "No visibility into federal contract spending patterns in Massachusetts that signal where SGC should focus bids"
]
txBox = s2.shapes.add_textbox(Inches(0.4), Inches(2.65), Inches(6.2), Inches(2.3))
tf = txBox.text_frame; tf.word_wrap = True; tf.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE; tf.clear()
for pt in pain_points:
    p = tf.add_paragraph(); p.text = "• " + pt; apply_font(p, 13, RGBColor(224, 230, 240)); p.space_after = Pt(8)

add_stat_box(s2, Inches(6.8), Inches(2.5), Inches(2.8), Inches(0.8), "1,621", "Vendors across unstructured lists")
add_stat_box(s2, Inches(6.8), Inches(3.4), Inches(2.8), Inches(0.8), "147", "Target companies with no unified view", RGBColor(12, 68, 124))
add_stat_box(s2, Inches(6.8), Inches(4.3), Inches(2.8), Inches(0.8), "$19.5B", "MA federal spend with no strategic map", RGBColor(15, 110, 86))

# ==========================================
# SLIDE 3: PROJECT OBJECTIVES
# ==========================================
s3 = prs.slides.add_slide(blank_slide_layout)
add_background(s3, COLOR_DARK_BLUE)
add_header(s3, "SLIDE 02", "Project Objectives", "Three guiding questions that shaped our analytical approach")
add_footer(s3)

add_card(s3, Inches(0.4), Inches(1.5), Inches(2.9), Inches(1.8), "01 — Market Density", 
         ["Which industry sectors (Healthcare, Tech, Finance) show the highest concentration of SGC Target companies for immediate consulting outreach?"], False)
add_card(s3, Inches(3.55), Inches(1.5), Inches(2.9), Inches(1.8), "02 — Vendor Optimization", 
         ["How can we visualize the vendor landscape to ensure a balanced mix of medical, logistical (VEH), and professional (PRF) service providers?"], False)
add_card(s3, Inches(6.7), Inches(1.5), Inches(2.9), Inches(1.8), "03 — Compliance & Diversity", 
         ["What is the distribution of SDO Commitment Percentages across the vendor pool, and how does this impact SGC's CSR reporting?"], False)

add_text(s3, "HIGH-LEVEL GOAL", Inches(0.4), Inches(3.5), Inches(4), Inches(0.3), 11, COLOR_PRIMARY_BLUE, bold=True)
add_highlight_box(s3, Inches(0.4), Inches(3.8), Inches(9.2), Inches(1.2), "Deliver a functional, data-driven dashboard that transforms static lists of thousands of companies and vendor contacts into an actionable Business Development Funnel — enabling SGC leadership to move from manual data lookup to strategic, educated decision-making regarding resource allocation and partnership strategy.")

# ==========================================
# SLIDE 4: DATA SOURCES
# ==========================================
s4 = prs.slides.add_slide(blank_slide_layout)
add_background(s4, COLOR_DARK_BLUE)
add_header(s4, "SLIDE 03", "Data Sources & Dataset Overview", "Primary datasets used to build the SGC Growth Intelligence Toolkit")
add_footer(s4)

add_card(s4, Inches(0.4), Inches(1.4), Inches(4.5), Inches(1.8), "Dataset 1 — Categorized Companies", 
         ["147 companies across 10 industry sectors", "Sectors: Healthcare, Finance, Tech, Energy, Education & more", 
          "Includes national, local, and SGC-targeted classifications", "58 priority accounts flagged for immediate outreach"], False)

add_card(s4, Inches(5.1), Inches(1.4), Inches(4.5), Inches(1.8), "Dataset 2 — Vendor Contact Details", 
         ["1,621 vendors across 15 service categories", "Service codes: MED, FAC, ITT, PRF, VEH & more", 
          "Includes contact names, roles, emails, phone numbers", "1,234 vendors with SDO commitment % data"], False)

add_card(s4, Inches(0.4), Inches(3.35), Inches(4.5), Inches(1.7), "Dataset 3 — USASpending.gov (Massachusetts)", 
         ["Federal contract & grant obligations by NAICS code", "Top 10 awarding agencies by spend in MA", 
          "HHS leads at $14.9B; Facilities Support tops NAICS at $19.5B"], False)

add_text(s4, "KEY DATASET STATS", Inches(5.1), Inches(3.35), Inches(4), Inches(0.3), 11, COLOR_PRIMARY_BLUE, bold=True)
stats_tags = ["147 Companies", "1,621 Vendors", "10 Sectors", "15 Service Codes", "58 SGC Targets", "Federal MA Spend", "SDO Percentage", "Contact Directory"]
for i, tag in enumerate(stats_tags):
    row, col = divmod(i, 4)
    tb = add_stat_box(s4, Inches(5.1) + col*Inches(1.1), Inches(3.7) + row*Inches(0.6), Inches(1.05), Inches(0.5), "", "", RGBColor(24, 60, 100))
    add_text(s4, tag, Inches(5.1) + col*Inches(1.1), Inches(3.7) + row*Inches(0.6), Inches(1.05), Inches(0.5), 9, RGBColor(91, 163, 232), bold=True, align=PP_ALIGN.CENTER, v_align=MSO_ANCHOR.MIDDLE, auto_size=True)

# ==========================================
# SLIDE 5: METHODOLOGY
# ==========================================
s5 = prs.slides.add_slide(blank_slide_layout)
add_background(s5, COLOR_DARK_BLUE)
add_header(s5, "SLIDE 04", "Methodology & System Architecture", "How raw data was transformed into a strategic decision-making toolkit")
add_footer(s5)

add_text(s5, "PIPELINE ARCHITECTURE", Inches(0.4), Inches(1.4), Inches(4), Inches(0.3), 11, COLOR_PRIMARY_BLUE, bold=True)
arch_steps = ["Raw Data Sources\n(Cos/Vendors/USASpending)", "Data Cleaning\n(Deduplication)", "Data Merging\n(Indexed by Service Code)", "Analysis & Viz\n(Chart.js/HTML)", "Interactive Dashboard\n(5-Tab Deployment)"]
for i, step in enumerate(arch_steps):
    bg = s5.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.4) + i*Inches(1.85), Inches(1.7), Inches(1.7), Inches(0.75))
    bg.fill.solid(); bg.fill.fore_color.rgb = RGBColor(20, 45, 80); bg.line.color.rgb = COLOR_PRIMARY_BLUE; bg.line.width = Pt(1.5)
    add_text(s5, step, Inches(0.4) + i*Inches(1.85), Inches(1.75), Inches(1.7), Inches(0.65), 10, RGBColor(200, 214, 232), align=PP_ALIGN.CENTER, v_align=MSO_ANCHOR.MIDDLE)
    if i < 4:
        add_text(s5, "→", Inches(0.4) + i*Inches(1.85) + Inches(1.7), Inches(1.8), Inches(0.15), Inches(0.5), 18, COLOR_PRIMARY_BLUE, bold=True, align=PP_ALIGN.CENTER)

add_text(s5, "ANALYTICAL METHODS", Inches(0.4), Inches(2.7), Inches(4), Inches(0.3), 11, COLOR_PRIMARY_BLUE, bold=True)
methods = ["Sector density analysis — company count by industry with SGC target flagging", "SDO tier segmentation — vendors bucketed into 5 diversity commitment levels", "National vs local presence mapping across all 10 sectors", "Federal spend alignment — NAICS obligations mapped to SGC service categories"]
txBox = s5.shapes.add_textbox(Inches(0.4), Inches(3.0), Inches(4.5), Inches(2.0))
tf = txBox.text_frame; tf.word_wrap = True; tf.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE; tf.clear()
for m in methods:
    p = tf.add_paragraph(); p.text = "• " + m; apply_font(p, 12.5, RGBColor(200, 214, 232)); p.space_after = Pt(8)

add_text(s5, "TOOLS & TECHNOLOGIES", Inches(5.2), Inches(2.7), Inches(4), Inches(0.3), 11, COLOR_PRIMARY_BLUE, bold=True)
tools = ["Python/Pandas", "Excel", "Chart.js", "HTML/CSS/JS", "USASpending.gov"]
for i, tag in enumerate(tools):
    row, col = divmod(i, 3)
    add_stat_box(s5, Inches(5.2) + col*Inches(1.4), Inches(3.1) + row*Inches(0.55), Inches(1.3), Inches(0.45), "", "", RGBColor(24, 60, 100))
    add_text(s5, tag, Inches(5.2) + col*Inches(1.4), Inches(3.1) + row*Inches(0.55), Inches(1.3), Inches(0.45), 10, RGBColor(91, 163, 232), bold=True, align=PP_ALIGN.CENTER, v_align=MSO_ANCHOR.MIDDLE)

add_highlight_box(s5, Inches(5.2), Inches(4.4), Inches(4.4), Inches(0.7), "All three datasets were merged into a unified relational structure indexed by service code, enabling cross-dataset filtering in the dashboard.")

# ==========================================
# SLIDE 6: VISUAL 1 & 2
# ==========================================
s6 = prs.slides.add_slide(blank_slide_layout)
add_background(s6, COLOR_LIGHT_GRAY)
add_header(s6, "SLIDE 05", "Visual Analysis — Sector Density & Company Presence", "Which sectors hold the most SGC target companies, and how does national vs local presence compare?", is_light=True)
add_footer(s6, is_light=True)

add_text(s6, "FIG 1 — COMPANY COUNT BY SECTOR", Inches(0.4), Inches(1.4), Inches(4.5), Inches(0.3), 11, COLOR_PRIMARY_BLUE, bold=True)
chart_data = CategoryChartData()
chart_data.categories = ["Healthcare", "Education", "Finance & Ins.", "Gov & NP", "Energy", "Construction", "Tech & Eng.", "Hospitality", "Publishing", "Misc"]
# To color individually, assign each category to its own series with 9 Nones!
c1_colors = ['#185FA5','#0F6E56','#3266AD','#BA7517','#534AB7','#639922','#993C1D','#888780','#D4537E','#1D9E75']
c1_vals = [21, 14, 27, 15, 20, 11, 24, 4, 5, 6]
for i in range(10):
    cd = [None]*10
    cd[i] = c1_vals[i]
    chart_data.add_series(f'S{i}', cd)

chart = s6.shapes.add_chart(XL_CHART_TYPE.COLUMN_STACKED, Inches(0.4), Inches(1.7), Inches(4.5), Inches(3.0), chart_data).chart
clean_chart(chart)
for i in range(10):
    chart.series[i].format.fill.solid()
    chart.series[i].format.fill.fore_color.rgb = hex_to_rgb(c1_colors[i])

add_text(s6, "Finding: Finance & Insurance leads with 27 companies; Technology & Engineering (24) and Healthcare (21) follow as top SGC opportunity sectors.", Inches(0.4), Inches(4.8), Inches(4.5), Inches(0.5), 11, COLOR_DARK_TEXT)

add_text(s6, "FIG 2 — NATIONAL VS LOCAL PRESENCE BY INDUSTRY", Inches(5.1), Inches(1.4), Inches(4.5), Inches(0.3), 11, COLOR_PRIMARY_BLUE, bold=True)
chart_data2 = CategoryChartData()
chart_data2.categories = ["Finance", "Tech", "Healthcare", "Energy", "Gov & NP", "Education", "Construction", "Misc", "Publishing", "Hospitality"][::-1]
chart_data2.add_series('National & Local', (29, 25, 21, 20, 25, 13, 11, 5, 5, 4)[::-1])
chart_data2.add_series('Local Only', (11, 4, 20, 9, 5, 7, 4, 2, 1, 1)[::-1])
chart_data2.add_series('SGC Target', (9, 9, 9, 8, 3, 9, 5, 2, 2, 2)[::-1])
chart2 = s6.shapes.add_chart(XL_CHART_TYPE.BAR_STACKED, Inches(5.1), Inches(1.7), Inches(4.5), Inches(3.0), chart_data2).chart
clean_chart(chart2)
chart2.has_legend = True; chart2.legend.position = XL_LEGEND_POSITION.BOTTOM
chart2.series[0].format.fill.solid(); chart2.series[0].format.fill.fore_color.rgb = hex_to_rgb('#00BCD4')
chart2.series[1].format.fill.solid(); chart2.series[1].format.fill.fore_color.rgb = hex_to_rgb('#7C4DFF')
chart2.series[2].format.fill.solid(); chart2.series[2].format.fill.fore_color.rgb = hex_to_rgb('#FF7043')

add_text(s6, "Finding: Finance & Insurance and Healthcare combined account for 18 of 58 SGC priority targets, with strong national and local representation in both.", Inches(5.1), Inches(4.8), Inches(4.5), Inches(0.5), 11, COLOR_DARK_TEXT)

# ==========================================
# SLIDE 7: VISUAL 3 & 4
# ==========================================
s7 = prs.slides.add_slide(blank_slide_layout)
add_background(s7, COLOR_LIGHT_GRAY)
add_header(s7, "SLIDE 06", "Visual Analysis — SDO Commitment & Diversity Compliance", "Distribution of supplier diversity commitments and which categories lead on CSR performance", is_light=True)
add_footer(s7, is_light=True)

add_text(s7, "FIG 3 — SDO COMMITMENT TIER DISTRIBUTION (1,234 VENDORS)", Inches(0.4), Inches(1.4), Inches(4.5), Inches(0.3), 11, COLOR_PRIMARY_BLUE, bold=True)
chart_data3 = CategoryChartData()
chart_data3.categories = ["0–5%", "5–10%", "10–20%", "20–50%", "50%+"]
chart_data3.add_series('Vendors', (676, 251, 224, 55, 28))
chart3 = s7.shapes.add_chart(XL_CHART_TYPE.DOUGHNUT, Inches(0.4), Inches(1.7), Inches(4.5), Inches(2.9), chart_data3).chart
clean_chart(chart3)
chart3.has_legend = True; chart3.legend.position = XL_LEGEND_POSITION.RIGHT
c3_colors = ['#B5D4F4','#378ADD','#185FA5','#0C447C','#042C53']
# Doughnut points can be colored individually!
for idx, pt in enumerate(chart3.series[0].points):
    pt.format.fill.solid()
    pt.format.fill.fore_color.rgb = hex_to_rgb(c3_colors[idx])

add_text(s7, "Finding: 54.8% of vendors commit less than 5% SDO, indicating significant room for SGC to champion diversity-first vendor selection in contract pitches.", Inches(0.4), Inches(4.7), Inches(4.5), Inches(0.5), 11, COLOR_DARK_TEXT)

add_text(s7, "FIG 4 — AVERAGE SDO % BY SERVICE CATEGORY", Inches(5.1), Inches(1.4), Inches(4.5), Inches(0.3), 11, COLOR_PRIMARY_BLUE, bold=True)
chart_data4 = CategoryChartData()
c4_cats = ["IT Equip.", "IT Services", "IT Telecom", "Facilities", "Vehicle", "Grocery", "Laundry", "Medical", "MRO", "Office", "Professional", "Print", "Staffing", "Waste"]
c4_vals = [5.8, 8.0, 5.1, 17.9, 3.1, 2.7, 3.7, 2.7, 4.8, 4.6, 7.3, 4.1, 13.2, 4.7]
chart_data4.categories = c4_cats
for i in range(len(c4_cats)):
    cd = [None]*len(c4_cats); cd[i] = c4_vals[i]; chart_data4.add_series(f'S{i}', cd)

chart4 = s7.shapes.add_chart(XL_CHART_TYPE.COLUMN_STACKED, Inches(5.1), Inches(1.7), Inches(4.5), Inches(3.0), chart_data4).chart
clean_chart(chart4)
for i, v in enumerate(c4_vals):
    c = '#042C53' if v>=15 else '#185FA5' if v>=8 else '#378ADD' if v>=5 else '#B5D4F4'
    chart4.series[i].format.fill.solid()
    chart4.series[i].format.fill.fore_color.rgb = hex_to_rgb(c)

add_text(s7, "Finding: Facilities (17.9%) and Staffing (13.2%) lead in diversity commitment — ideal segments for SGC's CSR-aligned contract pitches to government and non-profit clients.", Inches(5.1), Inches(4.8), Inches(4.5), Inches(0.5), 11, COLOR_DARK_TEXT)

# ==========================================
# SLIDE 8: VISUAL 5
# ==========================================
s8 = prs.slides.add_slide(blank_slide_layout)
add_background(s8, COLOR_LIGHT_GRAY)
add_header(s8, "SLIDE 07", "Visual Analysis — Massachusetts Federal Spending", "Where federal contract dollars flow in Massachusetts — mapping SGC's highest-value bidding opportunities", is_light=True)
add_footer(s8, is_light=True)

add_text(s8, "FIG 5 — TOP 10 NAICS SECTORS BY FEDERAL OBLIGATIONS IN MASSACHUSETTS (SOURCE: USASPENDING.GOV)", Inches(0.4), Inches(1.4), Inches(9.2), Inches(0.3), 11, COLOR_PRIMARY_BLUE, bold=True)
chart_data5 = CategoryChartData()
c5_cats = ["Facilities Support", "Bldg Construction", "Aircraft Mfg", "Health Insurance", "R&D Phys/Eng/Life", "Offices of Physicians", "Engineering Svc", "Computer Systems", "Pharma Mfg", "R&D Phys/Life Sci"][::-1]
c5_vals = [19.55, 17.92, 12.93, 11.05, 8.30, 7.68, 7.58, 7.02, 5.31, 4.96][::-1]
c5_colors = ['#185FA5','#0F6E56','#534AB7','#BA7517','#1D9E75','#639922','#3266AD','#993C1D','#888780','#D4537E'][::-1]
chart_data5.categories = c5_cats
for i in range(10):
    cd = [None]*10; cd[i] = c5_vals[i]; chart_data5.add_series(f'S{i}', cd)

chart5 = s8.shapes.add_chart(XL_CHART_TYPE.BAR_STACKED, Inches(0.4), Inches(1.7), Inches(9.2), Inches(2.7), chart_data5).chart
clean_chart(chart5)
for i in range(10):
    chart5.series[i].format.fill.solid()
    chart5.series[i].format.fill.fore_color.rgb = hex_to_rgb(c5_colors[i])

add_highlight_box(s8, Inches(0.4), Inches(4.5), Inches(3.0), Inches(0.7), "Facilities Support ($19.5B) — SGC's facilities vendor network directly maps to the single largest federal obligation category in MA.", is_light=True)
add_highlight_box(s8, Inches(3.5), Inches(4.5), Inches(3.0), Inches(0.7), "R&D & Engineering ($20.9B combined) — Professional services consulting is a natural SGC entry point into these high-spend categories.", is_light=True)
add_highlight_box(s8, Inches(6.6), Inches(4.5), Inches(3.0), Inches(0.7), "Healthcare & Pharma ($16.4B combined) — Aligns with SGC's 9 Healthcare/Biotech target accounts and HHS as the top awarding agency.", is_light=True)

# ==========================================
# SLIDE 9: KEY INSIGHTS
# ==========================================
s9 = prs.slides.add_slide(blank_slide_layout)
add_background(s9, COLOR_DARK_BLUE)
add_header(s9, "SLIDE 08", "Key Insights & Strategic Findings", "Data-driven takeaways that directly inform SGC's business development strategy")
add_footer(s9)

insights = [
    ("Finance & Healthcare dominate SGC's target landscape", "Finance & Insurance (27 companies, 9 targets) and Healthcare & Biotech (21 companies, 9 targets) represent the highest-density opportunity sectors for immediate outreach."),
    ("Facilities is the strongest SDO performer", "With an average SDO of 17.9% and Sustainable Facilities averaging 35.3%, this category offers SGC the most compelling diversity narrative for CSR-conscious clients."),
    ("54.8% of vendors fall below 5% SDO commitment", "A significant diversity gap exists across the vendor pool — SGC can position itself as a diversity-forward partner to address this gap in client supply chains."),
    ("MA federal spend signals high-value contract verticals", "Facilities Support ($19.5B), Building Construction ($17.9B), and Aircraft Manufacturing ($12.9B) are the top NAICS categories — all serviceable through SGC's current vendor network."),
    ("HHS & SSA represent 77.7% of top agency spend", "Healthcare-adjacent and social services consulting is the dominant federal investment category in MA — a direct strategic priority for SGC's government contract bidding."),
    ("IT sector holds highest-SDO individual vendors", "NEWCOM Wireless (51%), Digit Outsource (46%), and CenturyLink (36%) are standout diversity-committed IT vendors SGC should prioritize in tech contract proposals.")
]

for i, (title, desc) in enumerate(insights):
    row, col = divmod(i, 2)
    x = Inches(0.4) if col == 0 else Inches(5.1)
    y = Inches(1.5) + row * Inches(1.2)
    
    c = s9.shapes.add_shape(MSO_SHAPE.OVAL, x, y, Inches(0.4), Inches(0.4))
    c.fill.solid(); c.fill.fore_color.rgb = COLOR_PRIMARY_BLUE; c.line.fill.background()
    add_text(s9, str(i+1), x, y + Inches(0.05), Inches(0.4), Inches(0.3), 14, COLOR_WHITE, bold=True, align=PP_ALIGN.CENTER)
    
    add_text(s9, title, x + Inches(0.55), y, Inches(4), Inches(0.25), 13, COLOR_WHITE, bold=True)
    tb = add_text(s9, desc, x + Inches(0.55), y + Inches(0.25), Inches(4), Inches(0.8), 11, COLOR_MUTED_TEXT, auto_size=True)
    tb.text_frame.word_wrap = True

# ==========================================
# SLIDE 10: CONCLUSION
# ==========================================
s10 = prs.slides.add_slide(blank_slide_layout)
add_background(s10, COLOR_DARK_BLUE)
add_header(s10, "SLIDE 09", "Conclusion & Recommendations", "What SGC should do next with the Growth Intelligence Toolkit")
add_footer(s10)

add_text(s10, "STRATEGIC RECOMMENDATIONS", Inches(0.4), Inches(1.4), Inches(4.5), Inches(0.3), 11, COLOR_PRIMARY_BLUE, bold=True)
recs = [
    ("R1", "Prioritize Finance & Healthcare outreach", "Activate the 18 combined SGC targets across these two sectors immediately using the vendor contact funnel.", COLOR_PRIMARY_BLUE),
    ("R2", "Leverage SDO data in contract pitches", "Use Facilities and Staffing SDO leaders as proof points when pitching CSR-aligned contracts to clients.", RGBColor(15, 110, 86)),
    ("R3", "Target HHS & SSA-funded contracts in MA", "With $26.9B combined obligations, healthcare and social services federal contracts are SGC's highest-value bidding opportunity.", RGBColor(186, 117, 23)),
    ("R4", "Expand vendor SDO data coverage", "387 vendors currently lack SDO data — enriching these records will strengthen SGC's diversity reporting and vendor selection process.", RGBColor(83, 74, 183))
]

for i, (r_num, title, desc, colr) in enumerate(recs):
    y = Inches(1.8) + i * Inches(0.85)
    c = s10.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.4), y, Inches(0.4), Inches(0.4))
    c.fill.solid(); c.fill.fore_color.rgb = colr; c.line.fill.background()
    add_text(s10, r_num, Inches(0.4), y + Inches(0.05), Inches(0.4), Inches(0.3), 11, COLOR_WHITE, bold=True, align=PP_ALIGN.CENTER)
    
    add_text(s10, title, Inches(0.9), y, Inches(3.8), Inches(0.25), 12.5, COLOR_WHITE, bold=True)
    tb = add_text(s10, desc, Inches(0.9), y + Inches(0.25), Inches(3.8), Inches(0.55), 11, COLOR_MUTED_TEXT, auto_size=True)

add_text(s10, "CONCLUSION", Inches(5.1), Inches(1.4), Inches(4.5), Inches(0.3), 11, COLOR_PRIMARY_BLUE, bold=True)
add_highlight_box(s10, Inches(5.1), Inches(1.8), Inches(4.5), Inches(1.0), "The SGC Growth Intelligence Toolkit successfully transforms thousands of unstructured company and vendor records into a five-tab interactive dashboard that answers SGC's three core business questions on market density, vendor optimization, and diversity compliance.")
add_highlight_box(s10, Inches(5.1), Inches(2.9), Inches(4.5), Inches(0.9), "By integrating Massachusetts federal spending data from USASpending.gov, the toolkit goes beyond internal data to position SGC with a real-world contract bidding strategy grounded in $100B+ of publicly tracked federal obligations.")

add_text(s10, "NEXT STEPS", Inches(5.1), Inches(4.0), Inches(4.5), Inches(0.3), 11, COLOR_PRIMARY_BLUE, bold=True)
nstp = ["Deploy dashboard to SGC leadership team for live use", "Refresh vendor and company data quarterly", "Integrate Tableau/PowerBI version for enterprise-scale filtering"]
txBox = s10.shapes.add_textbox(Inches(5.1), Inches(4.3), Inches(4.5), Inches(1.0))
tf = txBox.text_frame; tf.word_wrap = True; tf.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE; tf.clear()
for n in nstp:
    p = tf.add_paragraph(); p.text = "• " + n; apply_font(p, 12, COLOR_WHITE); p.space_after = Pt(6)

# ==========================================
# SLIDE 11: REFERENCES
# ==========================================
s11 = prs.slides.add_slide(blank_slide_layout)
add_background(s11, COLOR_DARK_BLUE)
add_text(s11, "NORTHEASTERN UNIVERSITY · ALY6980 · WINTER 2026", Inches(1), Inches(0.8), Inches(8), Inches(0.3), 12, COLOR_PRIMARY_BLUE, bold=True, align=PP_ALIGN.CENTER)
add_text(s11, "Thank You", Inches(1), Inches(1.2), Inches(8), Inches(0.7), 48, COLOR_WHITE, bold=True, align=PP_ALIGN.CENTER)
add_text(s11, "We welcome any questions from SGC leadership and our faculty reviewers.", Inches(1), Inches(2.0), Inches(8), Inches(0.4), 16, COLOR_MUTED_TEXT, align=PP_ALIGN.CENTER)

refs = [
    "USASpending.gov — Federal Contract & Grant Obligations, Massachusetts, 2024",
    "Commonwealth of Massachusetts — Statewide Contracts Vendor Database (mass.gov)",
    "SGC Internal Dataset — Categorized Companies & Vendor Contact Details, Winter 2026",
    "NAICS Association — North American Industry Classification System Codes, 2024",
    "Northeastern University CPS — ALY6980 Capstone Project Guidelines, Winter 2026"
]

ref_box = s11.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(1.5), Inches(2.6), Inches(7), Inches(2.2))
ref_box.fill.solid(); ref_box.fill.fore_color.rgb = RGBColor(20, 35, 55); ref_box.line.color.rgb = COLOR_PRIMARY_BLUE; ref_box.line.width = Pt(1.5)

add_text(s11, "REFERENCES", Inches(1.8), Inches(2.8), Inches(6.4), Inches(0.3), 12, COLOR_PRIMARY_BLUE, bold=True)
txBox = s11.shapes.add_textbox(Inches(1.8), Inches(3.2), Inches(6.4), Inches(1.4))
tf = txBox.text_frame; tf.word_wrap = True; tf.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE; tf.clear()
for r in refs:
    p = tf.add_paragraph(); p.text = "• " + r; apply_font(p, 11, COLOR_MUTED_TEXT); p.space_after = Pt(4)

bar11 = s11.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), prs.slide_height - Inches(0.1), prs.slide_width, Inches(0.1))
bar11.fill.solid(); bar11.fill.fore_color.rgb = RGBColor(29, 158, 117); bar11.line.fill.background()

# Save PPT
prs.save('/Users/sumesh/Projects/Antigravity/Capstone/SGC_Growth_Intelligence_Toolkit.pptx')
print("Successfully generated beautifully formatted SGC_Growth_Intelligence_Toolkit.pptx")
